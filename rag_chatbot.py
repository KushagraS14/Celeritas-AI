import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFacePipeline
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, AutoModelForCausalLM, pipeline
import requests
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
import pandas as pd


# ------------------ File Extraction Functions ------------------

def get_txt_text(txt_files):
    text = ""
    for txt in txt_files:
        txt.seek(0)
        text += txt.read().decode("utf-8") + "\n\n"
    return text


def get_docx_text(docx_files):
    text = ""
    for doc in docx_files:
        doc.seek(0)
        document = docx.Document(doc)
        for para in document.paragraphs:
            text += para.text + "\n"
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
        text += "\n\n"
    return text


def get_pptx_text(pptx_files):
    text = ""
    for ppt in pptx_files:
        ppt.seek(0)
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        text += "\n\n"
    return text


def get_spreadsheet_text(spreadsheet_files):
    text = ""
    for f in spreadsheet_files:
        f.seek(0)
        try:
            filename = f.name.lower()
            if filename.endswith('.csv'):
                df = pd.read_csv(f)
            elif filename.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(f, engine='openpyxl' if filename.endswith('.xlsx') else None)
            else:
                continue
            buffer = []
            buffer.append(f"Sheet/Table: {f.name}")
            buffer.append("Columns: " + ", ".join(df.columns.tolist()))
            buffer.append("Sample Data (first 5 rows):")
            buffer.append(df.head().to_string(index=False))
            buffer.append("\n" + "-" * 50 + "\n")
            text += "\n".join(buffer)
        except Exception as e:
            st.warning(f"⚠️ Could not parse {f.name}: {str(e)}")
    return text


def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def get_all_text_from_files(uploaded_files):
    pdf_files, docx_files, txt_files, pptx_files, spreadsheet_files = [], [], [], [], []
    for f in uploaded_files:
        ext = f.name.lower().split('.')[-1]
        if ext == 'pdf': pdf_files.append(f)
        elif ext == 'docx': docx_files.append(f)
        elif ext == 'txt': txt_files.append(f)
        elif ext == 'pptx': pptx_files.append(f)
        elif ext in ['csv', 'xlsx', 'xls']: spreadsheet_files.append(f)
        
        else: st.warning(f"⚠️ Unsupported file type: {f.name}")

    raw_text = ""
    if pdf_files: raw_text += get_pdf_text(pdf_files) + "\n\n"
    if docx_files: raw_text += get_docx_text(docx_files) + "\n\n"
    if txt_files: raw_text += get_txt_text(txt_files) + "\n\n"
    if pptx_files: raw_text += get_pptx_text(pptx_files) + "\n\n"
    if spreadsheet_files: raw_text += get_spreadsheet_text(spreadsheet_files) + "\n\n"
    return raw_text


def get_url_text(urls):
    text = ""
    headers = {'User-Agent': 'Mozilla/5.0'}
    for url in urls:
        try:
            st.write(f"📥 Fetching: {url}")
            response = requests.get(url.strip(), headers=headers, timeout=10)
            response.raise_for_status()
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(response.text, 'html.parser')
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            page_text = ' '.join(soup.get_text(separator=' ', strip=True).split())
            text += page_text + "\n\n"
            st.write(f"✅ Fetched: {url}")
        except Exception as e:
            st.error(f"❌ Error fetching {url}: {str(e)}")
    return text


# ------------------ RAG and Chat Functions ------------------

def get_text_chunks(text):
    splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
    return splitter.split_text(text)


def get_vectorstore(chunks):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(texts=chunks, embedding=embeddings)


def get_conversation_chain(vectorstore, model_name="google/flan-t5-base", temperature=0.7, top_p=0.95, max_length=512):
    model_info = {
        "google/flan-t5-base": "seq2seq",
        "google/flan-t5-large": "seq2seq",
        "t5-small": "seq2seq",
        "gpt2": "causal",
        "microsoft/DialoGPT-medium": "causal"
    }
    model_type = model_info.get(model_name, "seq2seq")
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name) if model_type == "seq2seq" else AutoModelForCausalLM.from_pretrained(model_name)
    if model_type == "causal" and tokenizer.pad_token is None:
        tokenizer.pad_token = tokenizer.eos_token

    pipe = pipeline(
        "text2text-generation" if model_type == "seq2seq" else "text-generation",
        model=model,
        tokenizer=tokenizer,
        max_length=max_length,
        temperature=temperature,
        top_p=top_p,
        pad_token_id=tokenizer.pad_token_id
    )
    llm = HuggingFacePipeline(pipeline=pipe)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True, output_key='answer')
    return ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory, return_source_documents=True)


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    for i, message in enumerate(st.session_state.chat_history):
        template = user_template if i % 2 == 0 else bot_template
        st.write(template.replace("{{MSG}}", message.content), unsafe_allow_html=True)


# ------------------ Main App ------------------

def main():
    st.set_page_config(page_title="Multi-Source RAG Chatbot", page_icon="🤖", layout="wide")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    # ✅ Show sidebar only before instance creation
    if not st.session_state.get("chat_started", False):
        with st.sidebar:
            st.subheader("🏢 Instance Creation")
            instance_type = st.selectbox("Choose Domain", ["Healthcare", "Banking", "Broking"])
            if st.button("🚀 Create Instance", use_container_width=True):
                st.session_state.chat_started = True
                st.session_state.instance_type = instance_type
                st.rerun()

            st.subheader("📚 Document Sources")
            st.markdown("---")
            st.markdown("##### ⚙️ LLM Settings")

            model_options = [
                "google/flan-t5-base",
                "google/flan-t5-large",
                "t5-small",
                "gpt2",
                "microsoft/DialoGPT-medium"
            ]
            selected_model = st.selectbox("Choose an LLM:", model_options, index=0)
            temperature = st.slider("Temperature", 0.0, 1.0, 0.7, 0.01)
            top_p = st.slider("Top-p (Nucleus Sampling)", 0.0, 1.0, 0.95, 0.01)
            max_token_length = st.number_input("Max Token Length (Max 512)", 50, 512, 512, 1)

            tab1, tab2 = st.tabs(["📁 Files", "🌐 URLs"])

            with tab1:
                uploaded_files = st.file_uploader("Upload Documents", accept_multiple_files=True,
                                                  type=['pdf', 'docx', 'txt', 'pptx', 'csv', 'xlsx'])
                if st.button("Process Files", use_container_width=True):
                    if uploaded_files:
                        with st.spinner("🔄 Processing files..."):
                            raw_text = get_all_text_from_files(uploaded_files)
                            if raw_text.strip():
                                text_chunks = get_text_chunks(raw_text)
                                vectorstore = get_vectorstore(text_chunks)
                                st.session_state.conversation = get_conversation_chain(
                                    vectorstore, selected_model, temperature, top_p, max_token_length)
                                st.success("✅ Files processed!")
                            else:
                                st.error("❌ No text extracted.")
                    else:
                        st.warning("⚠️ Upload at least one file.")

            with tab2:
                url_input = st.text_area("Enter URLs (one per line)",
                                         placeholder="https://example.com\nhttps://another.com")
                if st.button("Process URLs", use_container_width=True):
                    if url_input.strip():
                        urls = [url.strip() for url in url_input.split('\n') if url.strip()]
                        with st.spinner("🔄 Fetching URLs..."):
                            raw_text = get_url_text(urls)
                            if raw_text.strip():
                                chunks = get_text_chunks(raw_text)
                                vectorstore = get_vectorstore(chunks)
                                st.session_state.conversation = get_conversation_chain(
                                    vectorstore, selected_model, temperature, top_p, max_token_length)
                                st.success("✅ URLs processed!")
                            else:
                                st.error("❌ No text found at URLs.")
                    else:
                        st.warning("⚠️ Enter at least one URL.")

    # ✅ If no instance started yet, show main interface
    if not st.session_state.get("chat_started", False):
        st.header("🤖 Multi-Source RAG Chatbot")
        st.markdown("_Ask questions about your uploaded documents or URLs!_")
        user_question = st.text_input("💬 Ask a question about your documents")
        if user_question:
            handle_userinput(user_question)
        return

    # ✅ After Create Instance → Sidebar is gone, clean chat window
    st.header(f"💬 {st.session_state.instance_type} Assistant")
    st.markdown(f"_You are now chatting in the **{st.session_state.instance_type}** domain._")

    if st.button("🔙 Back to Sidebar", use_container_width=True):
        st.session_state.chat_started = False
        st.session_state.conversation = None
        st.session_state.chat_history = None
        st.rerun()

    user_question = st.text_input("💬 Ask your question:")
    if user_question:
        if st.session_state.conversation is None:
            dummy_store = get_vectorstore(["This is a blank default context for testing."])
            st.session_state.conversation = get_conversation_chain(dummy_store)
        handle_userinput(user_question)


if __name__ == "__main__":
    main()
