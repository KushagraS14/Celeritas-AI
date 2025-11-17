import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFacePipeline
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM,  AutoModelForCausalLM, pipeline
import torch
import requests
from bs4 import BeautifulSoup
from langchain.docstore.document import Document
import docx
from pptx import Presentation
import pandas as pd
import os



# Function to extract text from .txt files
def get_txt_text(txt_files):
    text = ""
    for txt in txt_files:
        txt.seek(0)  # Reset file pointer
        text += txt.read().decode("utf-8") + "\n\n"
    return text


# Function to extract text from .docx files
def get_docx_text(docx_files):
    text = ""
    for doc in docx_files:
        doc.seek(0)
        document = docx.Document(doc)
        for para in document.paragraphs:
            text += para.text + "\n"
        # Also extract text from tables (optional but useful)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
        text += "\n\n"
    return text

# Function to extract text from .pptx files
def get_pptx_text(pptx_files):
    text = ""
    for ppt in pptx_files:
        ppt.seek(0)
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        text += "\n\n"
    return text


# Function to extract text from .csv and .xlsx files
def get_spreadsheet_text(spreadsheet_files):
    text = ""
    for f in spreadsheet_files:
        f.seek(0)
        try:
            # Determine file type by extension
            filename = f.name.lower()
            if filename.endswith('.csv'):
                df = pd.read_csv(f)
            elif filename.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(f, engine='openpyxl' if filename.endswith('.xlsx') else None)
            else:
                continue

            # Convert dataframe to a readable string format
            # Option: describe + sample rows
            buffer = []
            buffer.append(f"Sheet/Table: {f.name}")
            buffer.append("Columns: " + ", ".join(df.columns.tolist()))
            buffer.append("Sample Data (first 5 rows):")
            buffer.append(df.head().to_string(index=False))
            buffer.append("\n" + "-"*50 + "\n")
            text += "\n".join(buffer)
        except Exception as e:
            st.warning(f"⚠️ Could not parse {f.name}: {str(e)}")
    return text


# function to get texts from PDFs
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def get_all_text_from_files(uploaded_files):
    pdf_files = []
    docx_files = []
    txt_files = []
    pptx_files = []
    spreadsheet_files = []

    for f in uploaded_files:
        ext = f.name.lower().split('.')[-1]
        if ext == 'pdf':
            pdf_files.append(f)
        elif ext == 'docx':
            docx_files.append(f)
        elif ext == 'txt':
            txt_files.append(f)
        elif ext == 'pptx':
            pptx_files.append(f)
        elif ext in ['csv', 'xlsx', 'xls']:
            spreadsheet_files.append(f)
        else:
            st.warning(f"⚠️ Unsupported file type: {f.name}")

    raw_text = ""

    if pdf_files:
        raw_text += get_pdf_text(pdf_files) + "\n\n"
    if docx_files:
        raw_text += get_docx_text(docx_files) + "\n\n"
    if txt_files:
        raw_text += get_txt_text(txt_files) + "\n\n"
    if pptx_files:
        raw_text += get_pptx_text(pptx_files) + "\n\n"
    if spreadsheet_files:
        raw_text += get_spreadsheet_text(spreadsheet_files) + "\n\n"

    return raw_text


# function to fetch and extract text from URLs
def get_url_text(urls):
    text = ""
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
    }
    
    for url in urls:
        try:
            st.write(f"📥 Fetching: {url}")
            response = requests.get(url.strip(), headers=headers, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            
            page_text = soup.get_text(separator=' ', strip=True)
            page_text = ' '.join(page_text.split())
            text += page_text + "\n\n"
            st.write(f"✅ Successfully fetched: {url}")
            
        except Exception as e:
            st.error(f"❌ Error fetching {url}: {str(e)}")
    
    return text


# function to split the text into chunks
def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


# function to create vector store from the text chunks 
def get_vectorstore(text_chunks):
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore  


# function to create the conversation chain
def get_conversation_chain(vectorstore, selected_model_name, temperature, top_p, max_length):
    # model_name = "google/flan-t5-base"
    model_info = {
        "google/flan-t5-base": {"type": "seq2seq", "model": None, "tokenizer": None},
        "google/flan-t5-large": {"type": "seq2seq", "model": None, "tokenizer": None},
        "t5-small": {"type": "seq2seq", "model": None, "tokenizer": None},
        "gpt2": {"type": "causal", "model": None, "tokenizer": None},
        "microsoft/DialoGPT-medium": {"type": "causal", "model": None, "tokenizer": None},
        # Add more models here as needed
    }
    
    if selected_model_name not in model_info:
        st.error(f"Model {selected_model_name} is not configured.")
        return None

    model_type = model_info[selected_model_name]["type"]

    # Load tokenizer and model dynamically
    tokenizer = AutoTokenizer.from_pretrained(selected_model_name)
    if model_type == "seq2seq":
        model = AutoModelForSeq2SeqLM.from_pretrained(selected_model_name)
        task = "text2text-generation"
    elif model_type == "causal":
        model = AutoModelForCausalLM.from_pretrained(selected_model_name)
        task = "text-generation" # Different task for decoder-only models
        # Add padding token if it doesn't exist (common for GPT-style models)
        if tokenizer.pad_token is None:
            tokenizer.pad_token = tokenizer.eos_token

    # Create pipeline with user-defined parameters
    pipe = pipeline(
        task,
        model=model,
        tokenizer=tokenizer,
        max_length=max_length,  # Use user input
        temperature=temperature, # Use slider value
        top_p=top_p,             # Use slider value
        # For causal models, specify pad_token_id
        pad_token_id=tokenizer.pad_token_id if model_type == "causal" else None
    )

    # Wrap in LangChain
    llm = HuggingFacePipeline(pipeline=pipe)
    
    # Memory
    memory = ConversationBufferMemory(
        memory_key='chat_history', 
        return_messages=True,
        output_key='answer'  
    )
    
    # Create conversation chain
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory,
        return_source_documents=True
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)


def main():
    st.set_page_config(
        page_title="Multi-Source RAG Chatbot",
        page_icon="🤖",
        layout="wide"
    )
    
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    # Main header
    st.header("🤖 Multi-Source RAG Chatbot")
    # Parameter settings
    st.markdown("_Ask questions about your PDFs and websites!_")
    
    # Chat interface
    user_question = st.text_input("💬 Ask a question about your documents")
    if user_question:
        handle_userinput(user_question)

    # Sidebar
    with st.sidebar:
        st.subheader("📚 Document Sources")
        st.markdown("---")
        st.markdown("##### ⚙️ LLM Settings")

        model_options = [
            "google/flan-t5-base",
            "google/flan-t5-large",
            "t5-small",
            "gpt2",
            "microsoft/DialoGPT-medium"
        ]
        selected_model = st.selectbox(
            "Choose an LLM:",
            { "google/flan-t5-base",
            "google/flan-t5-large",
            "t5-small",
            "gpt2",
            "microsoft/DialoGPT-medium"},
            index=0,
            key="llm_selector"
        )
        temperature = st.slider(
            "Temperature",
            min_value=0.0,
            max_value=1.0,
            value=0.7,
            step=0.01,
            key="temp_slider"
        )
        top_p = st.slider(
            "Top-p (Nucleus Sampling)",
            min_value=0.0,
            max_value=1.0,
            value=0.95,
            step=0.01,
            key="top_p_slider"
        )
        max_token_length = st.number_input(
            "Max Token Length (Max 512)",
            min_value=50,
            max_value=512,
            value=512,
            step=1,
            key="max_token_input"
        )
        if max_token_length > 512:
            max_token_length = 512
            st.warning("Max token length capped at 512.")

        tab1, tab2 = st.tabs(["📁 Files", "🌐 URLs"])

        with tab1:
            st.markdown("##### Upload Documents")
            uploaded_files = st.file_uploader(
                "Supports: PDF, DOCX, TXT, PPTX, CSV, XLSX",
                accept_multiple_files=True,
                type=['pdf', 'docx', 'txt', 'pptx', 'csv', 'xlsx'],
                key="multi_uploader"
            )

            if st.button("Process Files", use_container_width=True):
                if uploaded_files:
                    with st.spinner("🔄 Processing files..."):
                        raw_text = get_all_text_from_files(uploaded_files)
                        if raw_text.strip():
                            text_chunks = get_text_chunks(raw_text)
                            vectorstore = get_vectorstore(text_chunks)
                            # Use the selected settings from the sidebar
                            st.session_state.conversation = get_conversation_chain(
                                vectorstore,
                                selected_model_name=selected_model,
                                temperature=temperature,
                                top_p=top_p,
                                max_length=max_token_length
                            )
                            # Optional: Store settings for display
                            if st.session_state.conversation is not None:
                                st.session_state.current_llm_name = selected_model
                                st.session_state.current_temp = temperature
                                st.session_state.current_top_p = top_p
                                st.session_state.current_max_tokens = max_token_length
                                st.success("✅ Files processed!")
                        else:
                            st.error("❌ No text extracted.")
                else:
                    st.warning("⚠️ Upload at least one file.")

        # URL Input Tab
        with tab2:
            st.markdown("##### Add Website URLs")
            url_input = st.text_area(
                "Enter URLs (one per line)",
                height=150,
                placeholder="https://example.com\nhttps://another-site.com",
                key="url_input"
            )
            
            if st.button("Process URLs", key="process_urls", use_container_width=True):
                if url_input.strip():
                    urls = [url.strip() for url in url_input.split('\n') if url.strip()]
                    
                    if urls:
                        with st.spinner("🔄 Fetching and processing URLs..."):
                            try:
                                # Get text from URLs
                                raw_text = get_url_text(urls)
                                
                                if raw_text.strip():
                                    # Get text chunks
                                    text_chunks = get_text_chunks(raw_text)
                                    st.success(f"✅ Extracted {len(text_chunks)} chunks from URLs")
                                    st.write(text_chunks)
                                    
                                    # Create vector store
                                    vectorstore = get_vectorstore(text_chunks)
                                    df = st.dataframe(vectorstore.index_to_docstore_id)
                                    st.write(df)

                                    
                                    # Create conversation chain
                                    # st.session_state.conversation = get_conversation_chain(vectorstore)
                                    st.session_state.conversation = get_conversation_chain(
                                        vectorstore,
                                        selected_model_name=selected_model,
                                        temperature=temperature,
                                        top_p=top_p,
                                        max_length=max_token_length
                                    )
                                     # Optional: Store settings for display
                                    if st.session_state.conversation is not None:
                                        st.session_state.current_llm_name = selected_model
                                        st.session_state.current_temp = temperature
                                        st.session_state.current_top_p = top_p
                                        st.session_state.current_max_tokens = max_token_length
                                        st.success("✅ Files processed!")
                                else:
                                    st.error("❌ No text could be extracted from URLs")
                            except Exception as e:
                                st.error(f"❌ Error processing URLs: {str(e)}")
                    else:
                        st.warning("⚠️ Please enter at least one valid URL")
                else:
                    st.warning("⚠️ Please enter at least one URL")
        
        # Combined Processing Option
        st.markdown("---")
        st.markdown("##### 🔗 Process Both")
        if st.button("Process Files + URLs Together", key="process_both", use_container_width=True):
            uploaded_files_combined = st.session_state.get('multi_uploader')
            url_input_combined = st.session_state.get('url_input', '')

            has_files = uploaded_files_combined and len(uploaded_files_combined) > 0
            has_urls = url_input_combined.strip()
            # pdf_docs_combined = st.session_state.get('pdf_uploader')
            # url_input_combined = st.session_state.get('url_input', '')
            
            # has_pdfs = pdf_docs_combined and len(pdf_docs_combined) > 0
            # has_urls = url_input_combined.strip()
            
            if has_files or has_urls:
                with st.spinner("🔄 Processing all sources..."):
                    try:
                        combined_text = ""
                        if has_files:
                            combined_text += get_all_text_from_files(uploaded_files_combined) + "\n\n"
                        if has_urls:
                            urls = [u.strip() for u in url_input_combined.split('\n') if u.strip()]
                            combined_text += get_url_text(urls)
                        
                        if combined_text.strip():
                            # Get text chunks
                            text_chunks = get_text_chunks(combined_text)
                            st.success(f"✅ Created {len(text_chunks)} total chunks")
                            st.write(text_chunks)
                            
                            # Create vector store
                            vectorstore = get_vectorstore(text_chunks)
                            
                            
                            # Create conversation chain
                            # st.session_state.conversation = get_conversation_chain(vectorstore)
                            st.session_state.conversation = get_conversation_chain(
                                vectorstore,
                                selected_model_name=selected_model,
                                temperature=temperature,
                                top_p=top_p,
                                max_length=max_token_length
                            )
                             # Optional: Store settings for display
                            if st.session_state.conversation is not None:
                                st.session_state.current_llm_name = selected_model
                                st.session_state.current_temp = temperature
                                st.session_state.current_top_p = top_p
                                st.session_state.current_max_tokens = max_token_length
                                st.success("✅ Files processed!")
                        else:
                            st.error("❌ No text could be extracted from any source")
                    except Exception as e:
                        st.error(f"❌ Error processing sources: {str(e)}")
            else:
                st.warning("⚠️ Please add PDFs or URLs first")
        
        # Info section
        st.markdown("---")
        st.markdown("##### ℹ️ How to use")
        st.markdown("""
        1. **Upload PDFs** or **Enter URLs** (or both!)
        2. Click the appropriate **Process** button
        3. Ask questions in the chat above
        4. The bot will answer based on your documents
        """)
        
    st.markdown("---")
st.subheader("🏢 Create Industry Instance")

industry = st.selectbox(
    "Choose Industry Domain",
    ["Healthcare", "Human Resources (HR)", "Banking", "Insurance"],
    key="industry_selector"
)

if st.button("🚀 Create Instance", use_container_width=True, key="create_instance_btn"):
    # Generate a unique URL per instance (could be different subpaths or query params)
    base_url = "http://localhost:8501"  # Change to your deployed Streamlit URL
    instance_url = f"{base_url}/?instance={industry.replace(' ', '_').lower()}"

    st.success(f"Launching {industry} Chat Instance...")
    
    # Open new tab using HTML+JS
    js = f"""
    <script>
        window.open('{instance_url}', '_blank').focus();
    </script>
    """
    st.markdown(js, unsafe_allow_html=True)

        # Clear conversation
if st.button("🗑️ Clear Conversation", use_container_width=True):
    st.session_state.conversation = None
    st.session_state.chat_history = None
    st.success("✅ Conversation cleared!")
    st.rerun()

                
if __name__ == "__main__":
    import urllib.parse

def main():
    st.set_page_config(page_title="Multi-Source RAG Chatbot", page_icon="🤖", layout="wide")

    # Read URL params
    query_params = st.experimental_get_query_params()
    instance = query_params.get("instance", ["default"])[0]
    instance_name = instance.replace("_", " ").title()

    st.write(css, unsafe_allow_html=True)

    # Display which instance this tab represents
    st.header(f"🤖 {instance_name} Chatbot Instance")

    main()